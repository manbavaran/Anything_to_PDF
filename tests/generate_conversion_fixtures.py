from pathlib import Path
import sys
import shutil
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from converter import PDFConverterLogic


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"
INPUTS = RESULTS / "inputs"
RESULTS.mkdir(exist_ok=True)
INPUTS.mkdir(exist_ok=True)


def make_image(path: Path, title: str, color: tuple[int, int, int]) -> None:
    image = Image.new("RGB", (1200, 700), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 1200, 110), fill=color)
    draw.text((40, 30), title, fill="white")
    draw.text((50, 160), "Anything to PDF 변환 테스트", fill="black")
    draw.text((50, 205), "이미지 내용: 색상, 제목, 설명 문장 포함", fill="black")
    draw.rectangle((50, 290, 1150, 580), outline=color, width=4)
    for row in range(3):
        for col in range(4):
            x, y = 80 + col * 260, 320 + row * 75
            draw.text((x, y), f"표 {row + 1}-{col + 1}: 테스트 값", fill="black")
    draw.text((50, 625), "[그림 1] 이미지 변환 및 병합 검증용 캡션", fill="black")
    image.save(path)


def make_docx(path: Path) -> None:
    doc = Document()
    doc.add_heading("Anything to PDF 문서 변환 테스트", 0)
    doc.add_paragraph("정품 Office 변환과 fallback 변환의 결과를 비교하기 위한 테스트 문서입니다.")
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    for cell, text in zip(table.rows[0].cells, ["항목", "설명", "값"]):
        cell.text = text
    for row in [("이미지", "색상 블록 포함", "정상"), ("표", "3열 표 포함", "정상"), ("캡션", "그림 설명 포함", "정상")]:
        cells = table.add_row().cells
        for cell, text in zip(cells, row):
            cell.text = text
    doc.add_paragraph("[표 1] DOCX 변환 테스트용 표 캡션")
    doc.add_picture(str(INPUTS / "01_test_image_red.png"), width=Inches(5.5))
    doc.add_paragraph("[그림 1] DOCX 안에 삽입된 테스트 이미지 캡션")
    doc.save(path)


def make_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Anything to PDF 프레젠테이션 테스트"
    textbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8), Inches(0.5))
    textbox.text_frame.text = "정품 Office 및 LibreOffice fallback 변환 비교용"
    table_shape = slide.shapes.add_table(4, 3, Inches(0.7), Inches(2.0), Inches(8), Inches(2.0))
    table = table_shape.table
    for col, text in enumerate(["구성", "내용", "상태"]):
        table.cell(0, col).text = text
    rows = [("이미지", "색상 블록", "포함"), ("표", "3열 데이터", "포함"), ("캡션", "그림 설명", "포함")]
    for row_index, row in enumerate(rows, 1):
        for col, text in enumerate(row):
            table.cell(row_index, col).text = text
    slide.shapes.add_picture(str(INPUTS / "02_test_image_blue.jpg"), Inches(0.8), Inches(4.4), width=Inches(3.2))
    caption = slide.shapes.add_textbox(Inches(4.2), Inches(5.2), Inches(5), Inches(0.5))
    caption.text_frame.text = "[그림 1] PPTX 변환 테스트 이미지 캡션"
    presentation.save(path)


def main() -> None:
    make_image(INPUTS / "01_test_image_red.png", "IMAGE TEST - RED", (180, 45, 55))
    make_image(INPUTS / "02_test_image_blue.jpg", "IMAGE TEST - BLUE", (35, 85, 180))
    # A second image is also a direct Pillow conversion fixture.
    make_image(INPUTS / "03_test_image_green.webp", "IMAGE TEST - GREEN", (35, 150, 85))
    make_docx(INPUTS / "04_test_table_image_caption.docx")
    make_pptx(INPUTS / "05_test_table_image_caption.pptx")

    logic = PDFConverterLogic()
    lines = ["Anything to PDF conversion test report", "", f"Results: {RESULTS}", ""]
    lines.append(f"DEFAULT_ENGINES | Word={logic._has_word()} | PowerPoint={logic._has_powerpoint()} | LibreOffice={logic._soffice_path or 'NOT_FOUND'} | Hancom={logic._has_hancom_hwp()}")
    lines.append("FALLBACK_NOTE | Office files use Microsoft Office when available; LibreOffice is selected only when the default Office engine is unavailable.")
    lines.append("")
    converted = []
    for source in sorted(INPUTS.iterdir()):
        engine = logic.engine_for(source)
        lines.append(f"SOURCE | {source.name} | ENGINE | {engine.label} | AVAILABLE | {engine.available}")
        if not engine.available:
            lines.append(f"SKIP   | {source.name} | reason={engine.note}")
            continue
        try:
            output = logic.convert_to_pdf(source)
            output_path = Path(output)
            saved_output = RESULTS / f"converted_{engine.name.replace(' ', '_')}_{source.stem}.pdf"
            shutil.copy2(output_path, saved_output)
            converted.append((source, saved_output))
            lines.append(f"CONVERT| {source.name} | OUTPUT | {saved_output.name} | ENGINE | {engine.name}")
        except Exception as exc:
            lines.append(f"ERROR  | {source.name} | {type(exc).__name__}: {exc}")

    direct_pdf = RESULTS / "06_direct_pdf_input.pdf"
    Image.new("RGB", (1000, 600), (245, 245, 245)).save(direct_pdf)
    converted.append((direct_pdf, direct_pdf))
    lines.append("SOURCE | 06_direct_pdf_input.pdf | ENGINE | PDF 직접 병합 | AVAILABLE | True")

    if converted:
        ordered = sorted(converted, key=lambda pair: pair[0].name)
        merged = RESULTS / "merged_01_name_order_Pillow_and_direct_PDF.pdf"
        logic.merge_pdfs([pdf for _, pdf in ordered], merged)
        lines.append("")
        lines.append("MERGE  | output=merged_01_name_order_Pillow_and_direct_PDF.pdf")
        lines.append("ORDER  | " + " -> ".join(source.name for source, _ in ordered))

        reversed_order = list(reversed(ordered))
        merged_reverse = RESULTS / "merged_02_reverse_order_Pillow_and_direct_PDF.pdf"
        logic.merge_pdfs([pdf for _, pdf in reversed_order], merged_reverse)
        lines.append("MERGE  | output=merged_02_reverse_order_Pillow_and_direct_PDF.pdf")
        lines.append("ORDER  | " + " -> ".join(source.name for source, _ in reversed_order))

    (RESULTS / "conversion_report.txt").write_text("\n".join(lines), encoding="utf-8")
    logic.cleanup_temps([])
    print("Created test fixtures and conversion report in", RESULTS)


if __name__ == "__main__":
    main()
